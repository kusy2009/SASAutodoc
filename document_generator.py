import re
import logging
import json
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.enum.table import WD_ALIGN_VERTICAL
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfbase import pdfmetrics
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
import os

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.DEBUG)

SECTION_ORDER = [
    "Overview",
    "Syntax",
    "Parameters",
    "Key Features and Functionalities",
    "Usage Examples",
    "Return Values",
    "Error Handling",
    "Version History",
    "Summary"
]

def create_presentation(content, output_path, preferences=None):
    """Creates a PowerPoint presentation from the documentation content."""
    try:
        prs = Presentation()

        # Set default preferences if none provided
        if not preferences:
            preferences = {
                'font_family': 'Arial',
                'font_size': 12,
                'heading_style': 'standard'
            }

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]  # Get subtitle placeholder

        # Always use macro name in title if available and not empty
        macro_name = content.get('macro_name')
        if macro_name and str(macro_name).strip():
            title.text = f"{macro_name} User Manual"
        else:
            title.text = "SAS Macro Documentation"
        subtitle.text = "SAS Macro Documentation"

        # Configure title formatting
        title_frame = title.text_frame
        title_frame.paragraphs[0].font.size = PptxPt(44)
        title_frame.paragraphs[0].font.name = preferences['font_family']

        # Content slides
        for section in SECTION_ORDER:
            if section not in content or not content[section]:
                continue

            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            title = slide.shapes.title
            title.text = section

            if section == "Parameters":
                if isinstance(content[section], dict):
                    rows = len(content[section]["table_rows"]) + 1
                    cols = len(content[section]["table_headers"])
                    shapes = slide.shapes
                    table = shapes.add_table(rows, cols, PptxInches(1), PptxInches(2),
                                             PptxInches(8), PptxInches(4)).table

                    # Add headers
                    for i, header in enumerate(content[section]["table_headers"]):
                        cell = table.cell(0, i)
                        cell.text = header

                    # Add data rows
                    for i, row_data in enumerate(content[section]["table_rows"], start=1):
                        for j, cell_data in enumerate(row_data):
                            cell = table.cell(i, j)
                            cell.text = str(cell_data)
            else:
                content_box = slide.placeholders[1]
                formatted_content = format_content(content[section])
                content_box.text = formatted_content

        prs.save(output_path)

    except Exception as e:
        logger.error(f"Error creating presentation: {str(e)}")
        raise

def create_manual(content, output_path, output_format='rtf', macro_name=None, preferences=None, rtf_mode=False):
    """Creates a document with formatted content."""
    try:
        if not isinstance(content, dict):
            logger.error(f"Invalid content provided to create_manual: {content}")
            raise ValueError("Invalid or empty content provided")

        # Store macro_name in content for use in other functions
        content['macro_name'] = macro_name

        if output_format == 'pdf':
            _create_pdf(content, output_path, macro_name, preferences)
        elif output_format == 'pptx':
            create_presentation(content, output_path, preferences)
        elif output_format == 'html':
             _create_html(content, output_path, macro_name, preferences)
        else:
            # For RTF, use the _create_rtf function with proper error handling
            try:
                _create_rtf(content, output_path, macro_name, preferences)
                # Verify file was created and has content
                if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                    raise Exception("RTF file was not created properly")
            except Exception as rtf_error:
                logger.error(f"RTF generation error: {str(rtf_error)}")
                raise Exception(f"Failed to generate RTF document: {str(rtf_error)}")

    except Exception as e:
        logger.error(f"Error in create_manual: {str(e)}")
        raise Exception(f"Failed to generate document: {str(e)}")

def _process_section_name(section):
    """Helper function to match sections with or without colons"""
    return section.rstrip(':').lower()

def _create_pdf(content, output_path, macro_name=None, preferences=None):
    """Creates a PDF document using ReportLab."""
    try:
        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            leftMargin=72,
            rightMargin=72,
            topMargin=72,
            bottomMargin=72
        )
        styles = getSampleStyleSheet()
        story = []

        # Apply preferences
        font_size = preferences.get('font_size', 12) if preferences else 12
        # Calculate spacing (matching RTF format)
        section_spacing = font_size * 0.8  # Reduced from 1.5 to 0.8 for less gap between sections
        paragraph_spacing = font_size * 0.5  # Reduced from 0.8 to 0.5 for better consistency

        # Create custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Title'],
            fontSize=font_size + 8,
            spaceAfter=section_spacing,
            leading=font_size * 1.1  # Keep reduced line spacing
        )

        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading1'],
            fontSize=font_size + 4,
            spaceBefore=section_spacing/2,  # Reduced space before headings
            spaceAfter=paragraph_spacing/2,  # Reduced space after headings
            leading=font_size * 1.1  # Keep reduced line spacing
        )

        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontSize=font_size,
            spaceBefore=paragraph_spacing/3,  # Further reduced spacing
            spaceAfter=paragraph_spacing/3,
            leading=font_size * 1.1  # Keep reduced line spacing
        )

        code_style = ParagraphStyle(
            'CustomCode',
            parent=styles['Code'],
            fontSize=font_size - 2,
            leftIndent=36,
            spaceBefore=paragraph_spacing/3,
            spaceAfter=paragraph_spacing/3,
            leading=font_size * 1.1  # Keep reduced line spacing
        )

        # Add title
        title = f"{macro_name} User Manual" if macro_name else "SAS Macro Documentation"
        story.append(Paragraph(title, title_style))

        # Add sections with reduced spacing
        for section in SECTION_ORDER:
            if section not in content or not content[section]:
                continue

            story.append(Paragraph(section, heading_style))

            if section == "Parameters":
                headers = content[section]["table_headers"]
                data = [headers] + content[section]["table_rows"]

                # Adjust column widths for better content fit
                col_widths = [
                    doc.width * 0.15,  # Parameter name 
                    doc.width * 0.12,  # Default 
                    doc.width * 0.50,  # Description 
                ]

                # Create table with adjusted style
                table = Table(data, colWidths=col_widths, repeatRows=1)
                table.setStyle(TableStyle([
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ('FONTSIZE', (0, 0), (-1, 0), font_size),  # Reduced header font size from font_size + 1
                    ('FONTSIZE', (0, 1), (-1, -1), font_size - 1),  # Content font size unchanged
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                    ('LEFTPADDING', (0, 0), (-1, -1), 4),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 4),
                    ('TOPPADDING', (0, 0), (-1, -1), 4),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
                    ('WORDWRAP', (0, 0), (-1, -1), True),
                    ('ROWHEIGHTS', (0, 0), (-1, -1), None)
                ]))
                story.append(table)
            else:
                formatted_content = format_content(content[section])
                paragraphs = formatted_content.split('\n')

                for para in paragraphs:
                    if para.strip():
                        if para.startswith('    '):  # Code block
                            story.append(Paragraph(para.strip(), code_style))
                        else:
                            story.append(Paragraph(para, normal_style))

            # Reduced space between sections (only a small gap)
            story.append(Spacer(1, section_spacing/2))

        doc.build(story)

    except Exception as e:
        logger.error(f"Error creating PDF: {str(e)}")
        raise

def _create_rtf(content, output_path, macro_name=None, preferences=None):
    """Creates an RTF document using python-docx."""
    try:
        doc = Document()
        logger.debug(f"Creating RTF document with content keys: {list(content.keys())}")

        # Apply preferences
        font_family = preferences.get('font_family', 'Arial') if preferences else 'Arial'
        font_size = preferences.get('font_size', 12) if preferences else 12
        heading_style = preferences.get('heading_style', 'standard') if preferences else 'standard'
        code_style_name = preferences.get('code_style', 'github') if preferences else 'github'

        # Create custom styles
        custom_code_style = 'CustomCode'
        if custom_code_style not in doc.styles:
            code_style = doc.styles.add_style(custom_code_style, 1)
            code_style.font.name = 'Courier New'
            code_style.font.size = Pt(font_size - 2)
            code_style.paragraph_format.space_before = Pt(3)
            code_style.paragraph_format.space_after = Pt(3)
            code_style.paragraph_format.left_indent = Inches(0.5)

        # Add title
        title = doc.add_heading(f"{macro_name} User Manual" if macro_name else "SAS Macro Documentation", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title.style.font.name = font_family
        title.style.font.size = Pt(font_size + 8)

        # Add each section
        for section in SECTION_ORDER:
            section_key = _process_section_name(section)
            matching_key = next((k for k in content.keys()
                               if _process_section_name(k) == section_key), None)

            if not matching_key or not content[matching_key]:
                logger.debug(f"Skipping empty or missing section: {section}")
                continue

            # Add section heading (with colon)
            heading = doc.add_heading(section, level=1)
            heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
            heading.style.font.name = font_family
            heading.style.font.size = Pt(font_size + 4)
            heading.paragraph_format.space_before = Pt(12)
            heading.paragraph_format.space_after = Pt(6)

            if section_key == "parameters":
                # Create parameters table
                table = doc.add_table(rows=1, cols=len(content[matching_key]["table_headers"]))
                table.style = 'Table Grid'
                table.autofit = False

                # Set header row
                header_cells = table.rows[0].cells
                for i, header in enumerate(content[matching_key]["table_headers"]):
                    header_cells[i].text = header
                    header_cells[i].paragraphs[0].runs[0].bold = True
                    header_cells[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                    header_cells[i].paragraphs[0].runs[0].font.name = font_family
                    header_cells[i].paragraphs[0].runs[0].font.size = Pt(font_size + 2)

                # Add data rows
                for row_data in content[matching_key]["table_rows"]:
                    cells = table.add_row().cells
                    for i, cell_data in enumerate(row_data):
                        cells[i].text = str(cell_data)
                        cells[i].vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                        cells[i].paragraphs[0].runs[0].font.name = font_family
                        cells[i].paragraphs[0].runs[0].font.size = Pt(font_size)

            else:
                # Format section content
                formatted_content = format_content(content[matching_key])
                paragraphs = formatted_content.split('\n')

                for para in paragraphs:
                    p = doc.add_paragraph()
                    if para.startswith('    '):  # Code block
                        p.style = doc.styles[custom_code_style]
                        run = p.add_run(para.strip())
                        run.font.name = 'Courier New'
                    else:
                        p.style = doc.styles['Normal']
                        p.text = para
                        for run in p.runs:
                            run.font.name = font_family
                            run.font.size = Pt(font_size)
                    # Adjust spacing for all paragraphs
                    p.paragraph_format.space_before = Pt(3)
                    p.paragraph_format.space_after = Pt(3)
                    p.paragraph_format.line_spacing = 1.0

        # Save document and verify
        doc.save(output_path)
        logger.debug(f"Successfully saved document to {output_path}")

        # Verify file was created and has content
        if not os.path.exists(output_path):
            raise Exception("RTF file was not created")
        if os.path.getsize(output_path) == 0:
            raise Exception("RTF file was created but is empty")

    except Exception as e:
        logger.error(f"Error creating RTF: {str(e)}")
        raise

def format_content(text):
    """Format content with proper line breaks and bullet points"""
    try:
        if isinstance(text, (list, dict)):
            # Handle Usage Examples list
            if isinstance(text, list):
                formatted_text = "\n\n".join(str(example).strip() for example in text)
                return formatted_text
            # Handle Key Features dictionary
            elif isinstance(text, dict) and "subsections" in text:
                main_section = text.get("main_section", "").strip()
                subsections = text["subsections"]

                formatted_text = main_section + "\n"
                for subsection in subsections:
                    # Format each feature as a bullet point without bold
                    formatted_text += f"- {subsection['title']}\n"
                    # Indent the description without extra spacing
                    formatted_text += f"  {subsection['description']}\n"
                return formatted_text.strip()

        if not isinstance(text, str):
            logger.warning(f"Unexpected text type in format_content: {type(text)}")
            text = str(text)

        lines = text.split('\n')
        formatted_lines = []
        in_example = False
        in_numbered_list = False
        in_code_block = False

        for line in lines:
            stripped_line = line.strip()

            # Skip empty lines at the start
            if not stripped_line and not formatted_lines:
                continue

            # Handle example code blocks
            if 'Example:' in stripped_line or 'Usage:' in stripped_line:
                if not in_example:
                    formatted_lines.append('\n' + stripped_line)
                    in_example = True
                continue

            # Handle code block starts
            if stripped_line.startswith('    '):
                if not in_code_block:
                    formatted_lines.append('    ' + stripped_line)
                    in_code_block = True
                else:
                    formatted_lines.append('    ' + stripped_line)

                continue
            else:
               in_code_block = False

            # Handle list items
            if stripped_line.startswith('- '):
                formatted_lines.append(stripped_line)
                continue

            # Handle numbered list
            if re.match(r'^\d+\.\s', stripped_line):
                formatted_lines.append(stripped_line)
                in_numbered_list = True
                continue
            else:
                in_numbered_list = False

            if in_example:
                formatted_lines.append(stripped_line)
                continue

            formatted_lines.append(stripped_line)

        return "\n".join(formatted_lines)

    except Exception as e:
        logger.error(f"Error in format_content: {str(e)}")
        return str(text)

def _create_html(content, output_path, macro_name=None, preferences=None):
    """Creates an HTML document for preview."""
    try:
        # Only include the documentation sections, not the feedback form
        preview_html = f"""
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>{macro_name} Documentation</title>
            <link href="https://stackpath.bootstrapcdn.com/bootstrap/4.5.2/css/bootstrap.min.css" rel="stylesheet">
            <style>
                /* Existing styles remain unchanged */
                .code-comments {{
                  background-color: #f8f9fa;
                  padding: 15px;
                  border-radius: 5px;
                  margin-bottom: 20px;
                  border: 1px solid #ddd;
                }}
                .comment-block {{
                    margin-bottom: 10px;
                    padding: 10px;
                    border-left: 3px solid #007bff;
                    background-color: #e9ecef;
                    border-radius: 3px;
                }}
                .comment-text {{
                   font-weight: bold;
                }}
                 .comment-line {{
                    font-size: 0.9em;
                    color: #6c757d;
                    margin-top: 5px;
                }}
                 .preview-content {{
                    padding: 20px;
                    font-family: Arial, sans-serif;
                    max-width: 900px;
                    margin: 20px auto;
                    background-color: white;
                    box-shadow: 0 0 10px rgba(0, 0, 0, 0.1);
                    border-radius: 8px;
                }}
                .table-responsive {{
                    overflow-x: auto;
                }}
                .table {{
                    width: 100%;
                    margin-bottom: 1rem;
                    color: #212529;
                }}
                .table th,
                .table td {{
                    padding: 0.75rem;
                    vertical-align: top;
                    border-top: 1px solid #dee2e6;
                    word-wrap: break-word;
                }}
                .table thead th {{
                    vertical-align: bottom;
                    border-bottom: 2px solid #dee2e6;
                }}
                .table-bordered {{
                  border: 1px solid #dee2e6;
                }}
                 .table-bordered th,
                 .table-bordered td {{
                   border: 1px solid #dee2e6;
                }}
                .text-center {{
                  text-align: center;
                }}
                 .mt-4 {{
                    margin-top: 24px;
                 }}
                 .mb-4 {{
                    margin-bottom: 24px;
                 }}
                 .mt-3 {{
                    margin-top: 16px;
                 }}
                 .mb-3 {{
                     margin-bottom: 16px;
                 }}
                 .text-wrap {{
                    white-space: pre-line; /* or use 'normal' if you prefer */
                 }}
            </style>
        </head>
        <body>
            <div class="preview-content">
                <h1 class="text-center mb-4">{macro_name} User Manual</h1>

                <h2>Overview</h2>
                <p>{content['Overview']}</p>

                <h2>Syntax</h2>
                <pre class="example-code">{content['Syntax']}</pre>

                <h2>Parameters</h2>
                <table class="table table-bordered">
                    <thead>
                        <tr>
                            {''.join(f'<th>{header}</th>' for header in content['Parameters']['table_headers'])}
                        </tr>
                    </thead>
                    <tbody>
                        {''.join(f'<tr>{"".join(f"<td>{cell}</td>" for cell in row)}</tr>' for row in content['Parameters']['table_rows'])}
                    </tbody>
                </table>

                <h2>Key Features and Functionalities</h2>
                <p>{content['Key Features and Functionalities']['main_section']}</p>
                {''.join(f'<h3>{subsection["title"]}</h3><p>{subsection["description"]}</p>' for subsection in content['Key Features and Functionalities']['subsections'])}

                <h2>Usage Examples</h2>
                {''.join(f'<pre class="example-code">{example}</pre>' for example in content['Usage Examples'])}

                <h2>Return Values</h2>
                <p>{content['Return Values']}</p>

                <h2>Error Handling</h2>
                <p>{content['Error Handling']}</p>

                <h2>Summary</h2>
                <p>{content['Summary']}</p>
            </div>
        </body>
        </html>
        """

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(preview_html)

    except Exception as e:
        logger.error(f"Error creating HTML preview: {str(e)}")
        raise
def create_preview_content(doc_content, macro_name=None, special_comments=None):
    """Creates an HTML preview of the documentation content."""
    try:
        if isinstance(doc_content, str):
            content = json.loads(doc_content)

        # Only include documentation sections
        html = f"""
        <div class="preview-content">
            <h1 class="text-center mb-4">{macro_name} User Manual</h1>
        """

        # Add each documentation section in order
        for section in SECTION_ORDER:
            if section not in content or not content[section]:
                continue

            html += f'<h2 class="mt-4 mb-3">{section}</h2>'

            if section == "Parameters":
                if isinstance(content[section], dict) and "table_headers" in content[section] and "table_rows" in content[section]:
                    html += '<div class="table-responsive"><table class="table table-bordered">'
                    html += '<thead><tr>'
                    for header in content[section]["table_headers"]:
                        html += f'<th>{header}</th>'
                    html += '</tr></thead><tbody>'

                    for row in content[section]["table_rows"]:
                        html += '<tr>'
                        for cell in row:
                            html += f'<td class="text-wrap">{str(cell)}</td>'
                        html += '</tr>'
                    html += '</tbody></table></div>'
                else:
                    html += f'<p class="mb-3">{content[section]}</p>'
            else:
                formatted_content = format_content(content[section])
                formatted_content = formatted_content.replace('**', '')
                formatted_content = formatted_content.replace('\n', '<br>')
                html += f'<p class="mb-3">{formatted_content}</p>'

        html += '</div>'
        return html
    except Exception as e:
        logger.error(f"Error creating preview content: {str(e)}")
        if isinstance(doc_content, str):
            logger.debug(f"Raw content: {doc_content}")
        raise